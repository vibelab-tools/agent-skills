#!/usr/bin/env python3
from __future__ import annotations

import argparse
import base64
import datetime as dt
import hashlib
import json
import mimetypes
import os
import re
import shutil
import sys
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any
from urllib.parse import unquote, urlparse


IMAGE_EXTENSIONS = {
    ".apng",
    ".avif",
    ".bmp",
    ".emf",
    ".gif",
    ".heic",
    ".heif",
    ".jpeg",
    ".jpg",
    ".png",
    ".svg",
    ".tif",
    ".tiff",
    ".webp",
    ".wmf",
}

DATA_URI_RE = re.compile(r"^data:([^;,]+)?(?:;charset=[^;,]+)?;base64,(.*)$", re.I | re.S)
MD_IMAGE_RE = re.compile(r"!\[([^\]]*)\]\(([^)\n]+)\)")


@dataclass
class Asset:
    id: str
    path: str
    relative_path: str
    kind: str
    source: str
    source_ref: str | None = None
    alt: str | None = None
    page: int | None = None
    slide: int | None = None
    sheet: str | None = None
    cell: str | None = None
    sha256: str | None = None
    linked: bool = False
    notes: list[str] = field(default_factory=list)


class AssetStore:
    def __init__(self, assets_dir: Path, absolute_links: bool, link_base: Path | None):
        self.assets_dir = assets_dir
        self.absolute_links = absolute_links
        self.link_base = link_base
        self.assets: list[Asset] = []
        self._hash_index: dict[str, Asset] = {}
        self._used_names: set[str] = set()

    def save_bytes(
        self,
        data: bytes,
        *,
        name: str,
        kind: str,
        source: str,
        source_ref: str | None = None,
        alt: str | None = None,
        page: int | None = None,
        slide: int | None = None,
        sheet: str | None = None,
        cell: str | None = None,
        notes: list[str] | None = None,
    ) -> Asset:
        digest = hashlib.sha256(data).hexdigest()
        existing = self._hash_index.get(digest)
        if existing is not None:
            return existing

        path = self._next_path(name, data)
        path.write_bytes(data)
        asset = Asset(
            id=f"asset-{len(self.assets) + 1:04d}",
            path=str(path.resolve()),
            relative_path=path.relative_to(self.assets_dir.parent).as_posix(),
            kind=kind,
            source=source,
            source_ref=source_ref,
            alt=alt,
            page=page,
            slide=slide,
            sheet=sheet,
            cell=cell,
            sha256=digest,
            notes=notes or [],
        )
        self.assets.append(asset)
        self._hash_index[digest] = asset
        return asset

    def copy_file(
        self,
        source_path: Path,
        *,
        name: str | None = None,
        kind: str,
        source: str,
        source_ref: str | None = None,
        alt: str | None = None,
        notes: list[str] | None = None,
    ) -> Asset:
        return self.save_bytes(
            source_path.read_bytes(),
            name=name or source_path.name,
            kind=kind,
            source=source,
            source_ref=source_ref,
            alt=alt,
            notes=notes,
        )

    def link_for(self, asset: Asset, alt: str | None = None) -> str:
        asset.linked = True
        target = asset.path
        if not self.absolute_links and self.link_base is not None:
            target = os.path.relpath(asset.path, self.link_base)
        target = Path(target).as_posix()
        escaped_alt = escape_alt(alt or asset.alt or asset.id)
        return f"![{escaped_alt}](<{target}>)"

    def _next_path(self, name: str, data: bytes) -> Path:
        safe = sanitize_filename(name)
        stem = Path(safe).stem or "asset"
        ext = Path(safe).suffix.lower()
        if ext not in IMAGE_EXTENSIONS:
            ext = guess_image_extension(data, ext or ".bin")

        candidate = f"{stem}{ext}"
        index = 2
        while candidate in self._used_names or (self.assets_dir / candidate).exists():
            candidate = f"{stem}-{index}{ext}"
            index += 1
        self._used_names.add(candidate)
        return self.assets_dir / candidate


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Convert a local file with MarkItDown and extract visual assets into local files.",
    )
    parser.add_argument("filename", help="Local file to convert.")
    parser.add_argument("-o", "--output", help="Write Markdown to this file. Defaults to stdout and a runtime run directory.")
    parser.add_argument("--assets-dir", help="Directory for extracted image assets.")
    parser.add_argument("--manifest", help="Path for the generated JSON manifest.")
    parser.add_argument("--runtime-root", default=str(Path.home() / ".vibelab-tools" / "agent-skills" / "markitdown"))
    parser.add_argument("-x", "--extension", help="MarkItDown extension hint, such as .pdf.")
    parser.add_argument("-m", "--mime-type", help="MarkItDown MIME type hint.")
    parser.add_argument("-c", "--charset", help="MarkItDown charset hint.")
    parser.add_argument("-p", "--use-plugins", action="store_true", help="Enable installed MarkItDown plugins.")
    parser.add_argument("--relative-links", action="store_true", help="Use links relative to the Markdown output directory.")
    parser.add_argument("--no-asset-section", action="store_true", help="Do not append the extracted asset index section.")
    parser.add_argument("--pdf-pages", choices=["render", "skip"], default="render", help="Render PDF pages as local PNG assets.")
    parser.add_argument("--pdf-dpi", type=int, default=144, help="DPI for rendered PDF page images.")
    parser.add_argument("--max-pdf-pages", type=int, default=0, help="Maximum PDF pages to render. 0 means no limit.")
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    input_path = Path(args.filename).expanduser().resolve()
    if not input_path.exists():
        print(f"Input file does not exist: {input_path}", file=sys.stderr)
        return 2
    if not input_path.is_file():
        print(f"Input path is not a file: {input_path}", file=sys.stderr)
        return 2

    output_path = Path(args.output).expanduser().resolve() if args.output else None
    run_dir = make_run_dir(args, input_path, output_path)
    assets_dir = Path(args.assets_dir).expanduser().resolve() if args.assets_dir else run_dir / "assets"
    manifest_path = Path(args.manifest).expanduser().resolve() if args.manifest else run_dir / "manifest.json"

    assets_dir.mkdir(parents=True, exist_ok=True)
    manifest_path.parent.mkdir(parents=True, exist_ok=True)
    if output_path is not None:
        output_path.parent.mkdir(parents=True, exist_ok=True)

    link_base = output_path.parent if output_path is not None else run_dir
    store = AssetStore(assets_dir, absolute_links=not args.relative_links, link_base=link_base)
    warnings: list[str] = []

    markdown = convert_with_markitdown(input_path, args)
    replacements: dict[str, list[Asset]] = {}
    extension = (args.extension or input_path.suffix).lower()

    try:
        if extension == ".pdf" and args.pdf_pages == "render":
            render_pdf_pages(input_path, store, args, warnings)
        elif extension == ".pptx":
            extract_pptx_images(input_path, store, replacements, warnings)
        elif extension == ".docx":
            extract_zip_member_images(input_path, "word/media/", "docx-image", store, warnings)
        elif extension == ".xlsx":
            extract_xlsx_images(input_path, store, warnings)
        elif extension in {".html", ".htm"}:
            extract_html_images(input_path, store, replacements, warnings)
        elif extension == ".epub":
            extract_zip_member_images(input_path, "", "epub-image", store, warnings, archive_source="epub")
            add_zip_replacement_keys(store.assets, replacements)
        elif extension == ".zip":
            extract_zip_member_images(input_path, "", "zip-image", store, warnings, archive_source="zip")
        elif extension in IMAGE_EXTENSIONS:
            store.copy_file(input_path, kind="standalone-image", source="image", source_ref=input_path.name, alt=input_path.stem)
    except Exception as exc:
        warnings.append(f"Asset extraction failed for {extension or input_path.suffix}: {exc}")

    markdown = rewrite_markdown_image_links(markdown, input_path.parent, store, replacements, warnings)

    if store.assets and not args.no_asset_section:
        markdown = append_asset_section(markdown, store.assets, manifest_path, store)

    manifest = {
        "input": str(input_path),
        "markdown": str(output_path or (run_dir / "document.md")),
        "assets_dir": str(assets_dir),
        "asset_count": len(store.assets),
        "assets": [asset.__dict__ for asset in store.assets],
        "warnings": warnings,
    }
    manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")

    final_output = output_path or (run_dir / "document.md")
    final_output.write_text(markdown, encoding="utf-8")
    if output_path is None:
        print(markdown)
    return 0


def make_run_dir(args: argparse.Namespace, input_path: Path, output_path: Path | None) -> Path:
    if args.assets_dir:
        return Path(args.assets_dir).expanduser().resolve().parent
    if output_path is not None:
        return output_path.parent / f"{output_path.stem}.assets"

    stamp = dt.datetime.now().strftime("%Y%m%dT%H%M%S")
    slug = sanitize_filename(input_path.stem)[:48] or "document"
    digest = hashlib.sha1(str(input_path).encode("utf-8")).hexdigest()[:8]
    return Path(args.runtime_root).expanduser().resolve() / "runs" / f"{stamp}-{slug}-{digest}"


def convert_with_markitdown(input_path: Path, args: argparse.Namespace) -> str:
    from markitdown import MarkItDown, StreamInfo

    stream_info = None
    if args.extension or args.mime_type or args.charset:
        stream_info = StreamInfo(
            extension=args.extension,
            mimetype=args.mime_type,
            charset=args.charset,
            filename=input_path.name,
        )

    md = MarkItDown(enable_plugins=args.use_plugins)
    result = md.convert(input_path, stream_info=stream_info, keep_data_uris=True)
    return result.markdown or result.text_content or ""


def render_pdf_pages(input_path: Path, store: AssetStore, args: argparse.Namespace, warnings: list[str]) -> None:
    try:
        import pypdfium2 as pdfium
    except ImportError as exc:
        warnings.append(f"Skipping PDF page rendering because pypdfium2 is unavailable: {exc}")
        return

    pdf = pdfium.PdfDocument(str(input_path))
    page_count = len(pdf)
    limit = page_count if args.max_pdf_pages <= 0 else min(page_count, args.max_pdf_pages)
    if limit < page_count:
        warnings.append(f"Rendered first {limit} of {page_count} PDF pages because --max-pdf-pages was set.")

    scale = max(args.pdf_dpi, 36) / 72.0
    for index in range(limit):
        page = pdf[index]
        bitmap = page.render(scale=scale)
        image = bitmap.to_pil()
        name = f"pdf-page-{index + 1:04d}.png"
        temp_path = store.assets_dir / name
        image.save(temp_path, format="PNG")
        data = temp_path.read_bytes()
        temp_path.unlink(missing_ok=True)
        store.save_bytes(
            data,
            name=name,
            kind="pdf-page-render",
            source="pdf",
            source_ref=f"page {index + 1}",
            alt=f"PDF page {index + 1}",
            page=index + 1,
        )
        close_quietly(image)
        close_quietly(bitmap)
        close_quietly(page)
    close_quietly(pdf)


def extract_pptx_images(
    input_path: Path,
    store: AssetStore,
    replacements: dict[str, list[Asset]],
    warnings: list[str],
) -> None:
    try:
        import pptx
    except ImportError as exc:
        warnings.append(f"Skipping PPTX image extraction because python-pptx is unavailable: {exc}")
        return

    presentation = pptx.Presentation(str(input_path))
    for slide_index, slide in enumerate(presentation.slides, start=1):
        image_index = 0
        for shape in sorted_shapes(slide.shapes):
            for picture in iter_picture_shapes(shape, pptx):
                image_index += 1
                image = picture.image
                image_ext = extension_from_name(image.filename) or extension_from_mime(image.content_type) or ".bin"
                asset = store.save_bytes(
                    image.blob,
                    name=f"slide-{slide_index:03d}-image-{image_index:03d}{image_ext}",
                    kind="embedded-image",
                    source="pptx",
                    source_ref=picture.name,
                    alt=f"Slide {slide_index} image {image_index}",
                    slide=slide_index,
                )
                placeholder = re.sub(r"\W", "", picture.name) + ".jpg"
                replacements.setdefault(placeholder, []).append(asset)


def sorted_shapes(shapes: Any) -> list[Any]:
    return sorted(
        shapes,
        key=lambda shape: (
            float("-inf") if not getattr(shape, "top", None) else shape.top,
            float("-inf") if not getattr(shape, "left", None) else shape.left,
        ),
    )


def iter_picture_shapes(shape: Any, pptx_module: Any) -> Any:
    if is_pptx_picture(shape, pptx_module):
        yield shape
    if shape.shape_type == pptx_module.enum.shapes.MSO_SHAPE_TYPE.GROUP:
        for subshape in sorted_shapes(shape.shapes):
            yield from iter_picture_shapes(subshape, pptx_module)


def is_pptx_picture(shape: Any, pptx_module: Any) -> bool:
    if shape.shape_type == pptx_module.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
        return True
    if shape.shape_type == pptx_module.enum.shapes.MSO_SHAPE_TYPE.PLACEHOLDER:
        return hasattr(shape, "image")
    return False


def extract_xlsx_images(input_path: Path, store: AssetStore, warnings: list[str]) -> None:
    try:
        import openpyxl
        from openpyxl.utils import get_column_letter
    except ImportError as exc:
        warnings.append(f"Skipping XLSX image extraction because openpyxl is unavailable: {exc}")
        return

    workbook = openpyxl.load_workbook(str(input_path), read_only=False, data_only=True)
    try:
        for worksheet in workbook.worksheets:
            for index, image in enumerate(getattr(worksheet, "_images", []), start=1):
                data = image._data()
                ext = extension_from_name(getattr(image, "path", "")) or f".{getattr(image, 'format', 'png')}"
                cell = None
                anchor = getattr(image, "anchor", None)
                marker = getattr(anchor, "_from", None)
                if marker is not None:
                    cell = f"{get_column_letter(marker.col + 1)}{marker.row + 1}"
                store.save_bytes(
                    data,
                    name=f"{sanitize_filename(worksheet.title)}-image-{index:03d}{ext}",
                    kind="embedded-image",
                    source="xlsx",
                    source_ref=f"{worksheet.title}!{cell or 'unknown'}",
                    alt=f"{worksheet.title} image {index}",
                    sheet=worksheet.title,
                    cell=cell,
                )
    finally:
        close_quietly(workbook)


def extract_html_images(
    input_path: Path,
    store: AssetStore,
    replacements: dict[str, list[Asset]],
    warnings: list[str],
) -> None:
    try:
        from bs4 import BeautifulSoup
    except ImportError as exc:
        warnings.append(f"Skipping HTML image extraction because beautifulsoup4 is unavailable: {exc}")
        return

    soup = BeautifulSoup(input_path.read_bytes(), "html.parser")
    for index, img in enumerate(soup.find_all("img"), start=1):
        src = img.get("src")
        if not src:
            continue
        alt = img.get("alt") or f"HTML image {index}"
        asset = save_uri_or_local_image(src, input_path.parent, store, f"html-image-{index:03d}", "html", alt)
        if asset is not None:
            replacements.setdefault(src, []).append(asset)
            replacements.setdefault(unquote(src), []).append(asset)


def save_uri_or_local_image(
    src: str,
    base_dir: Path,
    store: AssetStore,
    name: str,
    source: str,
    alt: str,
) -> Asset | None:
    data_uri = parse_data_uri(src)
    if data_uri is not None:
        mime_type, data = data_uri
        ext = extension_from_mime(mime_type) or guess_image_extension(data, ".bin")
        return store.save_bytes(data, name=f"{name}{ext}", kind="embedded-image", source=source, source_ref="data-uri", alt=alt)

    parsed = urlparse(src)
    if parsed.scheme in {"http", "https"}:
        return None
    if parsed.scheme == "file":
        candidate = Path(unquote(parsed.path))
    else:
        candidate = (base_dir / unquote(parsed.path)).resolve()
    if candidate.exists() and candidate.is_file() and candidate.suffix.lower() in IMAGE_EXTENSIONS:
        return store.copy_file(candidate, name=f"{name}{candidate.suffix.lower()}", kind="referenced-image", source=source, source_ref=src, alt=alt)
    return None


def extract_zip_member_images(
    input_path: Path,
    prefix: str,
    name_prefix: str,
    store: AssetStore,
    warnings: list[str],
    *,
    archive_source: str = "docx",
) -> None:
    try:
        with zipfile.ZipFile(input_path) as archive:
            members = [name for name in archive.namelist() if is_archive_image_member(name, prefix)]
            for index, member in enumerate(members, start=1):
                data = archive.read(member)
                ext = extension_from_name(member) or guess_image_extension(data, ".bin")
                store.save_bytes(
                    data,
                    name=f"{name_prefix}-{index:03d}{ext}",
                    kind="embedded-image",
                    source=archive_source,
                    source_ref=member,
                    alt=f"{archive_source.upper()} image {index}",
                )
    except zipfile.BadZipFile as exc:
        warnings.append(f"Skipping archive image extraction because the file is not a readable ZIP container: {exc}")


def is_archive_image_member(member: str, prefix: str) -> bool:
    normalized = member.replace("\\", "/")
    if prefix and not normalized.startswith(prefix):
        return False
    if normalized.endswith("/"):
        return False
    parts = normalized.split("/")
    if any(part in {"", ".", ".."} for part in parts):
        return False
    return Path(normalized).suffix.lower() in IMAGE_EXTENSIONS


def add_zip_replacement_keys(assets: list[Asset], replacements: dict[str, list[Asset]]) -> None:
    for asset in assets:
        if not asset.source_ref:
            continue
        replacements.setdefault(asset.source_ref, []).append(asset)
        replacements.setdefault(Path(asset.source_ref).name, []).append(asset)
        parts = asset.source_ref.split("/")
        for index in range(1, len(parts) - 1):
            replacements.setdefault("/".join(parts[index:]), []).append(asset)


def rewrite_markdown_image_links(
    markdown: str,
    base_dir: Path,
    store: AssetStore,
    replacements: dict[str, list[Asset]],
    warnings: list[str],
) -> str:
    def replace(match: re.Match[str]) -> str:
        alt = match.group(1)
        raw_target = match.group(2).strip()
        target = raw_target[1:-1] if raw_target.startswith("<") and raw_target.endswith(">") else raw_target

        replacement = pop_replacement(target, replacements)
        if replacement is not None:
            return store.link_for(replacement, alt=alt)

        data_uri = parse_data_uri(target)
        if data_uri is not None:
            mime_type, data = data_uri
            ext = extension_from_mime(mime_type) or guess_image_extension(data, ".bin")
            asset = store.save_bytes(data, name=f"markdown-image-{len(store.assets) + 1:03d}{ext}", kind="embedded-image", source="markdown", source_ref="data-uri", alt=alt)
            return store.link_for(asset, alt=alt)

        asset = save_uri_or_local_image(target, base_dir, store, f"markdown-image-{len(store.assets) + 1:03d}", "markdown", alt or "Markdown image")
        if asset is not None:
            return store.link_for(asset, alt=alt)

        return match.group(0)

    try:
        return MD_IMAGE_RE.sub(replace, markdown)
    except Exception as exc:
        warnings.append(f"Markdown image-link rewriting failed: {exc}")
        return markdown


def pop_replacement(target: str, replacements: dict[str, list[Asset]]) -> Asset | None:
    keys = [target, unquote(target), Path(target).name]
    for key in keys:
        values = replacements.get(key)
        if values:
            return values.pop(0)
    return None


def append_asset_section(markdown: str, assets: list[Asset], manifest_path: Path, store: AssetStore) -> str:
    lines = [
        "",
        "",
        "## Extracted Visual Assets",
        "",
        f"Manifest: `{manifest_path.resolve()}`",
        "",
    ]
    for asset in assets:
        label_parts = [asset.source]
        if asset.page is not None:
            label_parts.append(f"page {asset.page}")
        if asset.slide is not None:
            label_parts.append(f"slide {asset.slide}")
        if asset.sheet:
            label_parts.append(asset.sheet if not asset.cell else f"{asset.sheet}!{asset.cell}")
        label = " ".join(label_parts)
        lines.append(f"- {label}: {store.link_for(asset, alt=asset.alt)}")
    return markdown.rstrip() + "\n".join(lines) + "\n"


def parse_data_uri(value: str) -> tuple[str | None, bytes] | None:
    match = DATA_URI_RE.match(value.strip())
    if not match:
        return None
    mime_type = match.group(1)
    payload = re.sub(r"\s+", "", match.group(2))
    try:
        return mime_type, base64.b64decode(payload, validate=True)
    except Exception:
        return None


def sanitize_filename(name: str) -> str:
    name = Path(name).name
    name = re.sub(r"[^A-Za-z0-9._-]+", "-", name).strip(".-")
    return name or "asset"


def escape_alt(value: str) -> str:
    return re.sub(r"[\r\n\[\]]", " ", value).strip()


def extension_from_name(value: str | None) -> str | None:
    if not value:
        return None
    ext = Path(value).suffix.lower()
    return ext if ext else None


def extension_from_mime(mime_type: str | None) -> str | None:
    if not mime_type:
        return None
    if mime_type == "image/jpeg":
        return ".jpg"
    if mime_type == "image/svg+xml":
        return ".svg"
    return mimetypes.guess_extension(mime_type.split(";")[0].strip())


def guess_image_extension(data: bytes, fallback: str) -> str:
    if data.startswith(b"\x89PNG\r\n\x1a\n"):
        return ".png"
    if data.startswith(b"\xff\xd8\xff"):
        return ".jpg"
    if data.startswith((b"GIF87a", b"GIF89a")):
        return ".gif"
    if data.startswith(b"RIFF") and data[8:12] == b"WEBP":
        return ".webp"
    if data.startswith((b"II*\x00", b"MM\x00*")):
        return ".tiff"
    if data.startswith(b"BM"):
        return ".bmp"
    return fallback


def close_quietly(value: Any) -> None:
    close = getattr(value, "close", None)
    if close is None:
        return
    try:
        close()
    except Exception:
        pass


if __name__ == "__main__":
    raise SystemExit(main())
